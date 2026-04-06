#!/usr/bin/env python3
"""Generate a professional academic PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Constants ──────────────────────────────────────────────────────────────────
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Colors
BG_COLOR = RGBColor(0x1B, 0x1F, 0x2B)       # Deep navy/charcoal
GOLD = RGBColor(0xC9, 0xA9, 0x62)            # Warm gold
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
DIM_GRAY = RGBColor(0x99, 0x99, 0x99)
QUOTE_BG = RGBColor(0x14, 0x17, 0x20)        # Darker box for quotes
QUOTE_BORDER = RGBColor(0x8A, 0x76, 0x4A)    # Muted gold border
RULE_COLOR = GOLD

FONT_SERIF = "Times New Roman"

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]


# ── Helper Functions ───────────────────────────────────────────────────────────

def set_slide_bg(slide, color=BG_COLOR):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def set_text(tf, text, font_size=18, color=WHITE, bold=False, italic=False,
             alignment=PP_ALIGN.LEFT, font_name=FONT_SERIF, spacing_after=Pt(6)):
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    return p


def add_paragraph(tf, text, font_size=18, color=WHITE, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT, font_name=FONT_SERIF, spacing_after=Pt(6),
                  spacing_before=Pt(0)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    return p


def add_bullet(tf, text, font_size=18, color=WHITE, bold=False, italic=False,
               level=0, spacing_after=Pt(8)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = FONT_SERIF
    p.level = level
    p.space_after = spacing_after
    # Bullet character
    pPr = p._pPr
    if pPr is None:
        from pptx.oxml.ns import qn
        pPr = p._p.get_or_add_pPr()
    from pptx.oxml.ns import qn
    from lxml import etree
    buChar = etree.SubElement(pPr, qn('a:buChar'))
    buChar.set('char', '\u2022')
    buClr = etree.SubElement(pPr, qn('a:buClr'))
    srgb = etree.SubElement(buClr, qn('a:srgbClr'))
    srgb.set('val', 'C9A962')
    return p


def add_horizontal_rule(slide, left, top, width, color=RULE_COLOR, thickness=Pt(1.5)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, thickness
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_slide_number(slide, number):
    txBox = add_textbox(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.35))
    tf = txBox.text_frame
    set_text(tf, str(number), font_size=11, color=DIM_GRAY, alignment=PP_ALIGN.RIGHT)


def add_header(slide, text, top=Inches(0.4)):
    txBox = add_textbox(slide, Inches(0.8), top, Inches(11.7), Inches(0.7))
    tf = txBox.text_frame
    set_text(tf, text, font_size=32, color=GOLD, bold=True, alignment=PP_ALIGN.LEFT)
    # Rule below header
    add_horizontal_rule(slide, Inches(0.8), top + Inches(0.65), Inches(11.7), thickness=Pt(1))
    return txBox


def add_quote_box(slide, left, top, width, height, text, header_text=None):
    """Add a bordered, shaded textbox for source passages."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = QUOTE_BG
    shape.line.color.rgb = QUOTE_BORDER
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.margin_top = Inches(0.15)
    tf.margin_bottom = Inches(0.15)
    if header_text:
        set_text(tf, header_text, font_size=16, color=GOLD, bold=True,
                 spacing_after=Pt(8))
        add_paragraph(tf, text, font_size=15, color=LIGHT_GRAY, italic=True,
                      spacing_after=Pt(4))
    else:
        set_text(tf, text, font_size=15, color=LIGHT_GRAY, italic=True)
    return shape


# ── SLIDE 1: Title ────────────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide1)

# Title
txBox = add_textbox(slide1, Inches(1.5), Inches(1.8), Inches(10.3), Inches(2.0))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "Moral Corruption vs. Natural Necessity:", font_size=36, color=GOLD,
         bold=True, alignment=PP_ALIGN.CENTER, spacing_after=Pt(2))
add_paragraph(tf, "Two Models of World-Ending in Plato and Stoic Philosophy",
              font_size=36, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER,
              spacing_after=Pt(0))

# Decorative horizontal rule
add_horizontal_rule(slide1, Inches(3.0), Inches(3.9), Inches(7.3), thickness=Pt(2))

# Subtitle
txBox2 = add_textbox(slide1, Inches(1.5), Inches(4.2), Inches(10.3), Inches(0.6))
tf2 = txBox2.text_frame
set_text(tf2, "GESM: The End of the World", font_size=22, color=LIGHT_GRAY,
         alignment=PP_ALIGN.CENTER)

add_slide_number(slide1, 1)


# ── SLIDE 2: Introduction & Guiding Questions ─────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide2)
add_header(slide2, "The Central Problem")

txBox = add_textbox(slide2, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=10, spacing_after=Pt(0))

add_bullet(tf, "In Plato\u2019s Critias, Zeus personally intervenes to destroy Atlantis after observing its moral decline \u2014 destruction is a divine response to human corruption",
           font_size=20, spacing_after=Pt(20))

add_bullet(tf, "In Stoic cosmology, the world burns on a predetermined cosmic schedule, indifferent to human virtue or vice",
           font_size=20, spacing_after=Pt(20))

# Guiding question
txBox3 = add_textbox(slide2, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.8))
tf3 = txBox3.text_frame
tf3.word_wrap = True
set_text(tf3, "\u2018Why do worlds end \u2014 because of what we do, or in spite of it?\u2019",
         font_size=22, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide2, 2)


# ── SLIDE 3: Thesis Statement ─────────────────────────────────────────────────
slide3 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide3)
add_header(slide3, "Thesis")

# Top rule
add_horizontal_rule(slide3, Inches(1.2), Inches(1.8), Inches(10.9), thickness=Pt(1))

thesis_text = (
    "Plato\u2019s Critias frames the destruction of Atlantis as a direct consequence of moral "
    "corruption \u2014 the gradual dilution of divine nature in its rulers \u2014 whereas the Stoics "
    "present cosmic destruction as a morally neutral, mathematically inevitable conflagration "
    "built into the rational structure of the universe itself. This contrast reveals two "
    "fundamentally incompatible ancient answers to the question of why worlds end: one rooted "
    "in ethics, the other in physics."
)

txBox = add_textbox(slide3, Inches(1.2), Inches(2.1), Inches(10.9), Inches(4.0))
tf = txBox.text_frame
tf.word_wrap = True
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
set_text(tf, thesis_text, font_size=22, color=WHITE, alignment=PP_ALIGN.CENTER,
         spacing_after=Pt(0))

# Bottom rule
add_horizontal_rule(slide3, Inches(1.2), Inches(5.6), Inches(10.9), thickness=Pt(1))

add_slide_number(slide3, 3)


# ── SLIDE 4: Primary Source 1 — Plato's Critias ───────────────────────────────
slide4 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide4)
add_header(slide4, "Plato\u2019s Moral Model of Destruction")

# Left column — quote box
passage_text = (
    "\u2018When the divine portion began to grow faint as it was often blended with great "
    "quantities of mortality...they became disordered...inwardly they were filled with an "
    "unjust lust for possessions and power. But Zeus...could clearly see this state of "
    "affairs...and resolved to punish them.\u2019"
)
add_quote_box(slide4, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.8),
              passage_text, header_text="The Passage (Critias 121b\u2013c)")

# Right column
txBox = add_textbox(slide4, Inches(6.8), Inches(1.4), Inches(5.7), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "What This Illustrates", font_size=16, color=GOLD, bold=True,
         spacing_after=Pt(12))

add_bullet(tf, "Destruction is reactive \u2014 it responds to moral failure",
           font_size=18, spacing_after=Pt(16))
add_bullet(tf, "The mechanism (earthquake/flood) is natural but the decision is divine",
           font_size=18, spacing_after=Pt(16))
add_bullet(tf, "Human moral choice has direct cosmic consequences",
           font_size=18, spacing_after=Pt(16))

add_slide_number(slide4, 4)


# ── SLIDE 5: Analysis of Plato ────────────────────────────────────────────────
slide5 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide5)
add_header(slide5, "The Ethics of Atlantis\u2019 End")

txBox = add_textbox(slide5, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=10, spacing_after=Pt(0))

add_bullet(tf, "Atlantis begins virtuous, noble, and self-sufficient \u2014 destruction is not inevitable but contingent on behavior",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "The dilution of divine nature by mortality is a gradual moral process, not a physical one",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "This places Critias in a tradition where the world\u2019s end functions as judgment \u2014 a cosmic verdict on human civilization",
           font_size=20, spacing_after=Pt(16))

# Bottom note
txBox2 = add_textbox(slide5, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.7))
tf2 = txBox2.text_frame
tf2.word_wrap = True
set_text(tf2, "\u2018Destruction here is meaningful \u2014 it says something about who we are\u2019",
         font_size=18, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide5, 5)


# ── SLIDE 6: Primary Source 2 — Stoic Cosmology ───────────────────────────────
slide6 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide6)
add_header(slide6, "The Stoic Physical Model of Destruction")

passage_text2 = (
    "\u2018At certain fated times the entire world is subject to conflagration, and then is "
    "reconstituted afresh. The primary fire possesses the principles of all things and the "
    "causes of past, present, and future events. The nexus and succession of these is fate, "
    "knowledge, truth, and an inevitable and inescapable law.\u2019"
)
add_quote_box(slide6, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.8),
              passage_text2, header_text="The Passage (Long & Sedley 46G)")

# Right column
txBox = add_textbox(slide6, Inches(6.8), Inches(1.4), Inches(5.7), Inches(4.8))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "What This Illustrates", font_size=16, color=GOLD, bold=True,
         spacing_after=Pt(12))

add_bullet(tf, "Destruction is proactive \u2014 scheduled from the beginning of the cosmos",
           font_size=18, spacing_after=Pt(16))
add_bullet(tf, "The word \u2018fated\u2019 removes all moral causation",
           font_size=18, spacing_after=Pt(16))
add_bullet(tf, "Paradoxically, conflagration is the moment of maximum cosmic goodness and wisdom (46N)",
           font_size=18, spacing_after=Pt(16))

add_slide_number(slide6, 6)


# ── SLIDE 7: Analysis of Stoic Cosmology ──────────────────────────────────────
slide7 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide7)
add_header(slide7, "Destruction Without Judgment")

txBox = add_textbox(slide7, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.5))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=10, spacing_after=Pt(0))

add_bullet(tf, "Human virtue or vice is entirely irrelevant to when the conflagration occurs",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "After each conflagration, everything recurs identically \u2014 Socrates, Plato, every city (52C) \u2014 not a reset toward something better but an exact repetition",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "This is incompatible with the idea of destruction as punishment or moral consequence",
           font_size=20, spacing_after=Pt(16))

# Bottom note
txBox2 = add_textbox(slide7, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.7))
tf2 = txBox2.text_frame
tf2.word_wrap = True
set_text(tf2, "\u2018Destruction here is meaningless as judgment \u2014 it says nothing about who we are\u2019",
         font_size=18, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide7, 7)


# ── SLIDE 8: Secondary Source Discussion ──────────────────────────────────────
slide8 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide8)
add_header(slide8, "Secondary Source")

# Source citation
txBox = add_textbox(slide8, Inches(1.0), Inches(1.5), Inches(11.3), Inches(1.0))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "G\u00e1bor Betegh, \u2018Cosmological Ethics in the Timaeus and Early Stoicism\u2019,",
         font_size=18, color=LIGHT_GRAY, italic=True, alignment=PP_ALIGN.LEFT,
         spacing_after=Pt(2))
add_paragraph(tf, "Oxford Studies in Ancient Philosophy, 2003",
              font_size=18, color=LIGHT_GRAY, italic=True, spacing_after=Pt(4))

# Bullet points
txBox2 = add_textbox(slide8, Inches(1.0), Inches(2.8), Inches(11.3), Inches(3.5))
tf2 = txBox2.text_frame
tf2.word_wrap = True
set_text(tf2, "", font_size=10, spacing_after=Pt(0))

add_bullet(tf2, "Betegh argues that Plato\u2019s Timaeus establishes a cosmological framework in which the structure of the universe is itself ethical",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf2, "The Stoics inherit this framework but radically transform it \u2014 ethics and physics are unified differently",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf2, "This supports the paper\u2019s argument that the two models represent genuinely incompatible views on whether human moral behavior has cosmic weight",
           font_size=20, spacing_after=Pt(16))

add_slide_number(slide8, 8)


# ── SLIDE 9: The Deeper Stakes ────────────────────────────────────────────────
slide9 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide9)
add_header(slide9, "What Is Really at Stake?")

# Left column
left_box = slide9.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.4), Inches(5.5), Inches(4.5)
)
left_box.fill.solid()
left_box.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x30)
left_box.line.fill.background()

tf_left = left_box.text_frame
tf_left.word_wrap = True
tf_left.margin_left = Inches(0.25)
tf_left.margin_right = Inches(0.25)
tf_left.margin_top = Inches(0.2)
set_text(tf_left, "Plato\u2019s Model", font_size=20, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER, spacing_after=Pt(14))

add_bullet(tf_left, "Human choices have cosmic consequences", font_size=18, spacing_after=Pt(12))
add_bullet(tf_left, "The end of the world is a moral verdict", font_size=18, spacing_after=Pt(12))
add_bullet(tf_left, "Ethics matters at a universal scale", font_size=18, spacing_after=Pt(12))
add_bullet(tf_left, "Destruction is contingent \u2014 it could have been avoided", font_size=18, spacing_after=Pt(12))

# Right column
right_box = slide9.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(1.4), Inches(5.7), Inches(4.5)
)
right_box.fill.solid()
right_box.fill.fore_color.rgb = RGBColor(0x1F, 0x23, 0x30)
right_box.line.fill.background()

tf_right = right_box.text_frame
tf_right.word_wrap = True
tf_right.margin_left = Inches(0.25)
tf_right.margin_right = Inches(0.25)
tf_right.margin_top = Inches(0.2)
set_text(tf_right, "Stoic Model", font_size=20, color=GOLD, bold=True,
         alignment=PP_ALIGN.CENTER, spacing_after=Pt(14))

add_bullet(tf_right, "Human choices are cosmically irrelevant to destruction", font_size=18, spacing_after=Pt(12))
add_bullet(tf_right, "The end of the world is a physical event", font_size=18, spacing_after=Pt(12))
add_bullet(tf_right, "Ethics matters only for the individual soul", font_size=18, spacing_after=Pt(12))
add_bullet(tf_right, "Destruction is necessary \u2014 nothing could prevent it", font_size=18, spacing_after=Pt(12))

# Bottom question
txBox_bottom = add_textbox(slide9, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.7))
tf_b = txBox_bottom.text_frame
tf_b.word_wrap = True
set_text(tf_b, "\u2018Can the end of the world mean something, or does it just happen?\u2019",
         font_size=22, color=GOLD, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide9, 9)


# ── SLIDE 10: Conclusion ──────────────────────────────────────────────────────
slide10 = prs.slides.add_slide(blank_layout)
set_slide_bg(slide10)
add_header(slide10, "Conclusion")

txBox = add_textbox(slide10, Inches(1.0), Inches(1.4), Inches(11.3), Inches(4.2))
tf = txBox.text_frame
tf.word_wrap = True
set_text(tf, "", font_size=10, spacing_after=Pt(0))

add_bullet(tf, "These two models represent a genuine philosophical fork in ancient thought \u2014 either the cosmos is a moral arena or an indifferent rational machine",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "Cicero and Seneca attempt to bridge this gap, but the tension remains unresolved",
           font_size=20, spacing_after=Pt(16))
add_bullet(tf, "The question these texts raise is still urgent: does the universe care about human virtue, or are we simply along for the ride?",
           font_size=20, spacing_after=Pt(16))

# Final line
txBox2 = add_textbox(slide10, Inches(1.0), Inches(5.3), Inches(11.3), Inches(1.0))
tf2 = txBox2.text_frame
tf2.word_wrap = True
set_text(tf2, "\u2018The end of the world, in ancient thought, was never just physics \u2014 it was always also a question about us.\u2019",
         font_size=20, color=GOLD, italic=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide10, 10)


# ── Save ───────────────────────────────────────────────────────────────────────
output_path = "/home/user/Testing/Moral_Corruption_vs_Natural_Necessity.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
